from pptx import Presentation

def extract_from_pptx(file_path):
    prs = Presentation(file_path)
    documents = []
    for i, slide in enumerate(prs.slides):
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        if text.strip():
            documents.append({
                "text": text.strip(),
                "metadata": {"slide_number": i + 1}
            })
    return documents
